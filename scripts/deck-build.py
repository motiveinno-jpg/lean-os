# -*- coding: utf-8 -*-
"""
OwnerView 소개서(PPT) 빌더 — 2026-08-20 사장님 확정안
  · 26장(공용 24 + 투자자 2) · 16:9 1920x1080 기준
  · 참고 제안서 3건(와일리·STUDIO161·드림인사이트)의 문법: 풀블리드 표지 · CONTENTS ·
    챕터 구분 · 상단 러너 · 번호 매긴 기능 · 말풍선 콜아웃 · 3분할 카드 · 성과 수치 · 회사 정보 표
  · 사장님 지시: 설명 항목 1개 = 캡처 조각 1개(자기 섹션), 문구는 격식체, 이미지는 고화질(2배)
  · 데이터는 시연 회사 '오너뷰' — 모티브 실정보 노출 0

실행: python scripts/deck-build.py   →  dist/OwnerView_소개서_2026.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ── 색·글꼴 ──
BRAND = RGBColor(0x4F, 0x46, 0xE5)
DEEP  = RGBColor(0x31, 0x2E, 0x81)
NAVY  = RGBColor(0x1E, 0x1B, 0x4B)
INK   = RGBColor(0x16, 0x1A, 0x2B)
MUTED = RGBColor(0x58, 0x60, 0x74)
DIM   = RGBColor(0x98, 0xA0, 0xB5)
LINE  = RGBColor(0xE2, 0xE5, 0xEF)
SOFT  = RGBColor(0xF7, 0xF7, 0xFE)
AMBER = RGBColor(0xB4, 0x53, 0x09)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LAV   = RGBColor(0xC7, 0xD2, 0xFE)
GREEN = RGBColor(0x0F, 0x9D, 0x58)

FONT = "맑은 고딕"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
DOC = "2026 오너뷰 소개서"
CAP = "cap"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ── 유틸 ──
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    """runs: [(문자열, pt, bold, color)] 또는 [[(...)...]] 형태로 문단 나눔"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        for (t, size, bold, color) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
    return tb

def gradient_bg(s, c1, c2, angle=45):
    sp = rect(s, 0, 0, W, H, fill=c1)
    f = sp.fill
    f.gradient()
    f.gradient_angle = angle
    stops = f.gradient_stops
    stops[0].color.rgb = c1; stops[0].position = 0.0
    stops[1].color.rgb = c2; stops[1].position = 1.0
    return sp

def pic(s, path, x, y, maxw, maxh, border=True, shadow=True):
    """비율 유지하며 (x,y) 영역 안에 가운데 배치"""
    p = os.path.join(CAP, path)
    if not os.path.exists(p):
        ph = rect(s, x, y, maxw, maxh, fill=RGBColor(0xF1, 0xF2, 0xF8), line=LINE)
        text(s, x, y + maxh // 2, maxw, Inches(0.3), [(f"[{path} 없음]", 10, False, DIM)], PP_ALIGN.CENTER)
        return ph
    iw, ih = Image.open(p).size
    r = min(maxw / iw, maxh / ih)
    w_, h_ = int(iw * r), int(ih * r)
    px, py = x + (maxw - w_) // 2, y + (maxh - h_) // 2
    im = s.shapes.add_picture(p, px, py, w_, h_)
    if border:
        b = rect(s, px, py, w_, h_, fill=None, line=LINE)
        b.line.width = Pt(0.75)
    return im

def pic_fill(s, path, x, y, w, h, keep="left"):
    """영역을 꽉 채우고 넘치는 쪽만 잘라낸다(cover). 가로로 긴 표 조각이 카드 안에서
       우표처럼 작아지는 것을 막는다 — 글자가 읽히는 크기로 들어간다."""
    p = os.path.join(CAP, path)
    if not os.path.exists(p):
        rect(s, x, y, w, h, fill=RGBColor(0xF1, 0xF2, 0xF8), line=LINE)
        text(s, x, y + h // 2, w, Inches(0.3), [(f"[{path} 없음]", 10, False, DIM)], PP_ALIGN.CENTER)
        return None
    iw, ih = Image.open(p).size
    ar, tr = iw / ih, w / h
    im = s.shapes.add_picture(p, x, y, w, h)
    if ar > tr:                       # 가로가 넘친다 → 좌우를 자른다
        cut = 1 - tr / ar
        if keep == "left":
            im.crop_right = cut
        else:
            im.crop_left = cut / 2; im.crop_right = cut / 2
    elif ar < tr:                     # 세로가 넘친다 → 위/아래를 자른다(위쪽을 남긴다)
        cut = 1 - ar / tr
        im.crop_bottom = cut
    b = rect(s, x, y, w, h, fill=None, line=LINE)
    b.line.width = Pt(0.75)
    return im

def runner(s, cat):
    text(s, Inches(0.62), Inches(0.34), Inches(6), Inches(0.24), [(cat, 9.5, True, BRAND)])
    text(s, Inches(6.7), Inches(0.34), Inches(6), Inches(0.24), [(DOC, 9, False, DIM)], PP_ALIGN.RIGHT)
    ln = rect(s, Inches(0.62), Inches(0.62), Inches(12.09), Emu(9525), fill=LINE)
    return ln

def footer(s, n):
    text(s, Inches(0.62), Inches(6.95), Inches(3), Inches(0.24), [("OwnerView", 8.5, True, BRAND)])
    text(s, Inches(9.7), Inches(6.95), Inches(3), Inches(0.24), [(f"{n:02d} / 26", 8.5, False, DIM)], PP_ALIGN.RIGHT)

def head(s, title_runs, lead=None, y=Inches(0.95)):
    text(s, Inches(0.62), y, Inches(12.1), Inches(0.5), title_runs)
    if lead:
        text(s, Inches(0.62), y + Inches(0.52), Inches(12.1), Inches(0.3), [(lead, 11, False, MUTED)])

def numtag(n):
    return (f"{n}. ", 20, True, BRAND)


# ══════════════════════════════════════════════════════════
# 01 표지
# ══════════════════════════════════════════════════════════
s = slide()
gradient_bg(s, NAVY, BRAND, angle=45)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill=LAV)
text(s, Inches(0.85), Inches(0.72), Inches(8), Inches(0.3),
     [("2026 · COMPANY OPERATION OS", 11, True, LAV)])
text(s, Inches(0.85), Inches(2.35), Inches(7.6), Inches(2.1),
     [[("은행·카드·세금·계약·직원 —", 30, True, WHITE)],
      [("회사 운영, 한 화면에서 끝냅니다.", 30, True, LAV)]], spacing=1.25)
text(s, Inches(0.85), Inches(4.35), Inches(7.4), Inches(0.4),
     [("흩어진 회사 운영을 하나로 모으는 운영 OS · www.owner-view.com", 12, False, RGBColor(0xD8, 0xDC, 0xF5))])
lg = rect(s, Inches(0.85), Inches(5.55), Inches(0.42), Inches(0.42), fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.85), Inches(5.63), Inches(0.42), Inches(0.3), [("O", 16, True, BRAND)], PP_ALIGN.CENTER)
text(s, Inches(1.4), Inches(5.62), Inches(4), Inches(0.4), [("OwnerView", 20, True, WHITE)])
pic(s, "01-dashboard-full.png", Inches(7.35), Inches(1.55), Inches(6.6), Inches(4.4), border=False)

# ══════════════════════════════════════════════════════════
# 02 CONTENTS
# ══════════════════════════════════════════════════════════
s = slide(); runner(s, "CONTENTS"); footer(s, 2)
head(s, [("목차", 22, True, INK)], "회사 소개부터 도입 비용까지 네 갈래로 나누어 말씀드립니다.")
items = [("Ⅰ", "회사 소개", "누가 만들고, 무엇을 만드는 회사인지"),
         ("Ⅱ", "무엇을 해결합니까", "회사의 숫자가 늦게 도착하는 이유와 해답"),
         ("Ⅲ", "주요 기능", "화면 하나하나가 없애는 일 (01~08)"),
         ("Ⅳ", "도입과 비용", "도입 절차 · 요금제 · 문의")]
for i, (no, t, d) in enumerate(items):
    x = Inches(0.62) + (i % 2) * Inches(6.2)
    y = Inches(2.35) + (i // 2) * Inches(1.75)
    text(s, x, y, Inches(0.6), Inches(0.5), [(no, 20, True, BRAND)])
    text(s, x + Inches(0.72), y, Inches(5.2), Inches(0.4), [(t, 15, True, INK)])
    text(s, x + Inches(0.72), y + Inches(0.42), Inches(5.2), Inches(0.4), [(d, 10.5, False, DIM)])
    rect(s, x, y + Inches(1.0), Inches(5.7), Emu(9525), fill=LINE)

# ══════════════════════════════════════════════════════════
# 챕터 구분 슬라이드
# ══════════════════════════════════════════════════════════
def chapter(no_txt, title, sub, page):
    s = slide()
    gradient_bg(s, INK, DEEP, angle=30)
    rect(s, Inches(0), Inches(0), Inches(0.16), H, fill=BRAND)
    text(s, Inches(1.1), Inches(2.75), Inches(8), Inches(0.35), [(no_txt, 12, True, LAV)])
    text(s, Inches(1.1), Inches(3.15), Inches(9), Inches(0.9), [(title, 32, True, WHITE)])
    text(s, Inches(1.1), Inches(4.15), Inches(9), Inches(0.4), [(sub, 12, False, RGBColor(0xC7, 0xCB, 0xE0))])
    text(s, Inches(9.7), Inches(6.95), Inches(3), Inches(0.24),
         [(f"{page:02d} / 26", 8.5, False, RGBColor(0x8A, 0x90, 0xB5))], PP_ALIGN.RIGHT)
    return s

chapter("CHAPTER Ⅰ", "회사 소개", "오너뷰를 만드는 회사와 서비스를 소개합니다.", 3)

# ══════════════════════════════════════════════════════════
# 04 회사 정보
# ══════════════════════════════════════════════════════════
s = slide(); runner(s, "회사 소개"); footer(s, 4)
head(s, [("회사 정보", 22, True, INK)], "오너뷰는 주식회사 모티브이노베이션이 만들고 운영합니다.")
rows = [("회사명", "주식회사 모티브이노베이션 (Motive Innovation Co., Ltd.)"),
        ("대표자", "채희웅"),
        ("설립", "2021년 9월"),
        ("소재지", "경기도 화성시 동탄첨단산업1로 27, A동 2514~2515호"),
        ("업태 · 종목", "전문, 과학 및 기술서비스업 · 광고대행 / 정보통신업 · 소프트웨어 개발"),
        ("서비스", "OwnerView — 회사 운영 통합 소프트웨어 (SaaS)"),
        ("홈페이지", "www.owner-view.com"),
        ("문의", "creative@mo-tive.com")]
y0 = Inches(2.25)
for i, (k, v) in enumerate(rows):
    y = y0 + Inches(0.5) * i
    text(s, Inches(0.72), y, Inches(2.2), Inches(0.35), [(k, 11, False, DIM)])
    text(s, Inches(3.0), y, Inches(9.5), Inches(0.35), [(v, 11.5, True, INK)])
    rect(s, Inches(0.72), y + Inches(0.4), Inches(11.9), Emu(9525), fill=LINE)

# ══════════════════════════════════════════════════════════
# 05 오너뷰 한 문장 + 6영역
# ══════════════════════════════════════════════════════════
s = slide(); runner(s, "회사 소개"); footer(s, 5)
head(s, [("오너뷰는 ", 22, True, INK), ("회사 운영 전체를 하나로 모으는 ", 22, True, BRAND), ("소프트웨어입니다.", 22, True, INK)],
     "은행·카드·세금·계약·직원·프로젝트가 한 데이터 위에서 돌아갑니다.")
areas = [("파이낸스", "통장 · 카드 · 거래처 · 정기 지출", "은행과 카드사에서 자료가 자동으로 들어옵니다."),
         ("기장 · 세무", "수집 · 전표 · 세금계산서 · 부가세", "분류와 전표를 제안하고, 확정은 사람이 합니다."),
         ("분석", "대시보드 · 경영흐름 · 손익 · 원장", "지금 상태와 30일 뒤를 함께 봅니다."),
         ("워크", "프로젝트 · 일정 · 게시판 · 파일", "일마다 맞는 화면으로 기록이 쌓입니다."),
         ("인사", "근태 · 연차 · 급여 · 근로계약", "직원은 마이페이지 하나로 끝냅니다."),
         ("전자계약 · 결재", "서명 · 직인 · 결재 허브", "종이 없이 문서가 흐름을 따라 갑니다.")]
for i, (t, sub, d) in enumerate(areas):
    x = Inches(0.62) + (i % 3) * Inches(4.05)
    y = Inches(2.5) + (i // 3) * Inches(2.05)
    rect(s, x, y, Inches(3.85), Inches(1.8), fill=SOFT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(s, x + Inches(0.25), y + Inches(0.22), Inches(3.4), Inches(0.3), [(t, 13.5, True, BRAND)])
    text(s, x + Inches(0.25), y + Inches(0.62), Inches(3.4), Inches(0.3), [(sub, 9.5, False, DIM)])
    text(s, x + Inches(0.25), y + Inches(1.0), Inches(3.4), Inches(0.6), [(d, 10.5, False, MUTED)], spacing=1.25)

# ══════════════════════════════════════════════════════════
# 06 챕터 Ⅱ
# ══════════════════════════════════════════════════════════
chapter("CHAPTER Ⅱ", "무엇을 해결합니까", "회사의 숫자가 늦게 도착하는 이유부터 말씀드립니다.", 6)

# ══════════════════════════════════════════════════════════
# 07 문제
# ══════════════════════════════════════════════════════════
s = slide(); runner(s, "무엇을 해결합니까"); footer(s, 7)
text(s, Inches(0.62), Inches(1.0), Inches(12), Inches(0.6),
     [("“", 26, True, BRAND), ("지금, 회사에 돈이 얼마나 있습니까?", 26, True, INK), ("”", 26, True, BRAND)])
text(s, Inches(0.62), Inches(1.75), Inches(12), Inches(0.3),
     [("이 질문에 바로 답하기 어려운 데에는 이유가 있습니다.", 11.5, False, MUTED)])
chaos = [("은행앱 3개", "잔액은 앱마다"), ("엑셀 장부", "지난주 상태"), ("세무사 문자", "월말에 한 번"), ("근태앱 따로", "급여는 수기"),
         ("종이 계약서", "만기를 놓칩니다"), ("견적서 폴더", "버전이 제각각"), ("단톡방 결재", "기록이 사라집니다"), ("흩어진 숫자", "즉답이 어렵습니다")]
for i, (t, d) in enumerate(chaos):
    x = Inches(0.62) + (i % 4) * Inches(3.06)
    y = Inches(2.65) + (i // 4) * Inches(1.35)
    rect(s, x, y, Inches(2.86), Inches(1.15), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    text(s, x, y + Inches(0.24), Inches(2.86), Inches(0.3), [(t, 12.5, True, INK)], PP_ALIGN.CENTER)
    text(s, x, y + Inches(0.62), Inches(2.86), Inches(0.3), [(d, 10, False, MUTED)], PP_ALIGN.CENTER)
rect(s, Inches(0.62), Inches(5.55), Inches(12.1), Inches(0.75), fill=RGBColor(0xFD, 0xF3, 0xEC), line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
text(s, Inches(0.62), Inches(5.75), Inches(12.1), Inches(0.4),
     [("도구 7개 · 손 입력 · 월말에야 도착하는 숫자 — 결정이 감(感)에 기대게 됩니다", 12, True, AMBER)], PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 08 해답 3단계
# ══════════════════════════════════════════════════════════
s = slide(); runner(s, "무엇을 해결합니까"); footer(s, 8)
head(s, [("오너뷰는 ", 22, True, INK), ("모으고 · 제안하고 · 사람이 확정합니다", 22, True, BRAND)],
     "자동화의 목적은 사람을 대신하는 것이 아니라, 판단할 준비를 끝내 두는 것입니다.")
steps = [("① 자동 수집", "은행 · 카드 · 국세청(세금계산서) 자료가 스스로 들어옵니다.", False),
         ("② AI 제안", "분류 · 매칭 · 오늘 챙길 것을 출처와 함께 제안합니다.\n(AI 추천 / 배운 규칙 / 국세청 / 장부 대조)", False),
         ("③ 사람이 확정", "발송 · 이체 · 확정 버튼은 항상 사람의 손에 있습니다.\n확정 버튼은 화면에 하나뿐입니다.", True)]
for i, (t, d, hl) in enumerate(steps):
    x = Inches(0.62) + i * Inches(4.28)
    rect(s, x, Inches(2.6), Inches(3.85), Inches(2.5),
         fill=SOFT if hl else WHITE, line=BRAND if hl else LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(s, x + Inches(0.28), Inches(2.95), Inches(3.3), Inches(0.4), [(t, 15, True, BRAND if hl else INK)])
    text(s, x + Inches(0.28), Inches(3.55), Inches(3.3), Inches(1.3), [(d, 10.5, False, MUTED)], spacing=1.3)
    if i < 2:
        text(s, x + Inches(3.95), Inches(3.6), Inches(0.35), Inches(0.4), [("→", 16, True, DIM)], PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 09 챕터 Ⅲ
# ══════════════════════════════════════════════════════════
chapter("CHAPTER Ⅲ", "주요 기능", "화면 하나하나가 어떤 일을 없애는지 말씀드립니다.", 9)


# ── 항목별 섹션 장표(★ 사장님 요청 레이아웃) ──
def feature_sections(page, no, title_txt, accent_txt, lead, sections, cat="주요 기능"):
    """sections: [(소제목, 이미지파일, 설명, 강조어)] — 항목마다 제 이미지와 제 섹션"""
    s = slide(); runner(s, cat); footer(s, page)
    text(s, Inches(0.62), Inches(0.95), Inches(12.1), Inches(0.45),
         [(f"{no}. ", 20, True, BRAND), (title_txt, 20, True, INK), (accent_txt, 20, True, BRAND)])
    text(s, Inches(0.62), Inches(1.52), Inches(12.1), Inches(0.3), [(lead, 11, False, MUTED)])
    n = len(sections)
    gap = Inches(0.22)
    total = Inches(12.1)
    cw = int((total - gap * (n - 1)) / n)
    for i, (t, img, d, hl) in enumerate(sections):
        x = Inches(0.62) + i * (cw + gap)
        y = Inches(2.05)
        card_h = Inches(4.55)
        rect(s, x, y, cw, card_h, fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
        # 머리
        rect(s, x, y, cw, Inches(0.72), fill=SOFT, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        rect(s, x, y + Inches(0.5), cw, Inches(0.22), fill=SOFT)
        text(s, x + Inches(0.2), y + Inches(0.08), cw - Inches(0.4), Inches(0.22), [(f"{i + 1:02d}", 9, True, BRAND)])
        text(s, x + Inches(0.2), y + Inches(0.32), cw - Inches(0.4), Inches(0.3), [(t, 12.5, True, INK)])
        rect(s, x, y + Inches(0.72), cw, Emu(9525), fill=LINE)
        # 이미지 — 카드 안 넓은 영역을 쓰되, 설명이 들어갈 자리는 남긴다
        img_top = y + Inches(0.88)
        img_h = Inches(2.42)
        pic(s, img, x + Inches(0.16), img_top, cw - Inches(0.32), img_h)
        # 설명 — 이미지 아래 고정 위치(카드마다 같은 줄에서 시작해 시선이 흔들리지 않는다)
        runs = [(d, 10, False, MUTED)]
        if hl:
            runs = [(d, 10, False, MUTED), (" " + hl, 10, True, AMBER)]
        text(s, x + Inches(0.2), y + Inches(3.52), cw - Inches(0.4), Inches(0.95), runs, spacing=1.3)
    return s

def feature_rows(page, no, title_txt, accent_txt, lead, sections, cat="주요 기능"):
    """가로로 긴 캡처(표·차트)에 맞는 형식 — 한 행 = 이미지 + 설명. 항목마다 제 이미지와 제 섹션.
       (카드형은 세로 조각에, 이 형식은 가로 조각에 쓴다 — 사장님 확정 목업의 두 가지 형식)"""
    s = slide(); runner(s, cat); footer(s, page)
    text(s, Inches(0.62), Inches(0.95), Inches(12.1), Inches(0.45),
         [(f"{no}. ", 20, True, BRAND), (title_txt, 20, True, INK), (accent_txt, 20, True, BRAND)])
    text(s, Inches(0.62), Inches(1.52), Inches(12.1), Inches(0.3), [(lead, 11, False, MUTED)])
    n = len(sections)
    top = Inches(2.05)
    gap = Inches(0.18)
    row_h = int((Inches(4.5) - gap * (n - 1)) / n)
    for i, (t, img, d, hl) in enumerate(sections):
        y = top + i * (row_h + gap)
        rect(s, Inches(0.62), y, Inches(12.1), row_h, fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        # 번호 · 소제목 (왼쪽)
        text(s, Inches(0.9), y + Inches(0.16), Inches(2.3), Inches(0.22), [(f"{i + 1:02d}", 8.5, True, BRAND)])
        text(s, Inches(0.9), y + Inches(0.38), Inches(2.7), Inches(0.32), [(t, 12.5, True, INK)])
        runs = [(d, 9.5, False, MUTED)]
        if hl:
            runs = [(d, 9.5, False, MUTED), (" " + hl, 9.5, True, AMBER)]
        text(s, Inches(0.9), y + Inches(0.76), Inches(2.75), row_h - Inches(0.85), runs, spacing=1.22)
        # 이미지 (오른쪽 넓게) — 가로로 긴 표는 채우고(왼쪽 기준), 세로가 있는 조각은 비율을 지킨다
        ip = os.path.join(CAP, img)
        ar = (lambda t: t[0] / t[1])(Image.open(ip).size) if os.path.exists(ip) else 9
        if ar >= 3.2:
            pic_fill(s, img, Inches(3.85), y + Inches(0.14), Inches(8.65), row_h - Inches(0.28), keep="left")
        else:
            pic(s, img, Inches(3.85), y + Inches(0.14), Inches(8.65), row_h - Inches(0.28))
    return s

# 10 · 01 통장·카드 자동 연동
feature_rows(10, "01", "통장·카드가 ", "스스로 들어옵니다",
                 "은행과 카드사에서 자료를 받아오는 일부터 장부에 앉히는 일까지, 사람 손을 거치지 않습니다.",
                 [("전 계좌를 한 화면에", "x/bank-accounts.png",
                   "여러 은행 앱을 오갈 필요가 없습니다. 계좌별 잔액과 이번 달 변화가 한 줄로 섭니다.", "매일 확인 0초"),
                  ("법인카드도 함께", "x/cards-list.png",
                   "카드사 승인 내역이 자동으로 들어오고, 카드별 결제일과 한도를 함께 관리합니다.", ""),
                  ("분류까지 자동", "x/collect-head.png",
                   "적요를 읽어 계정과목과 거래처를 제안합니다. 분류 완료율이 화면에 그대로 보입니다.", "손 입력 0건")])

# 11 · 02 기장 자동화
feature_rows(11, "02", "장부는 ", "채워진 채로 시작합니다",
                 "수집된 자료가 전표 후보까지 만들어 둡니다. 확정은 담당자가 한 번에 처리합니다.",
                 [("자료 수집", "x/collect-head.png",
                   "통장·카드·세금계산서·현금영수증을 채널별로 모으고, 무엇이 몇 건 들어왔는지 보여 줍니다.", ""),
                  ("세금계산서", "x/tax-rows.png",
                   "발행과 수취 내역이 한 표에 모입니다. 매입매출전표로 바로 이어집니다.", ""),
                  ("전표 확정", "x/voucher-rows.png",
                   "매입매출전표가 제안된 상태로 놓입니다. 사람이 확인하고 확정하면 장부가 됩니다.", "확정은 사람이")])

# 12 · 03 회사 신호 (콜아웃형)
s = slide(); runner(s, "주요 기능"); footer(s, 12)
text(s, Inches(0.62), Inches(0.95), Inches(12.1), Inches(0.45),
     [("03. ", 20, True, BRAND), ("회사 상태를 ", 20, True, INK), ("여섯 칸으로 봅니다", 20, True, BRAND)])
text(s, Inches(0.62), Inches(1.52), Inches(12.1), Inches(0.3),
     [("로그인하면 가장 먼저 만나는 줄입니다. 통장 잔액부터 세금 기한까지 한 줄에 있습니다.", 11, False, MUTED)])
pic(s, "x/dash-signals.png", Inches(0.62), Inches(2.2), Inches(12.1), Inches(1.15))
notes = [("통장 잔액", "지금 쓸 수 있는 돈"), ("30일 뒤 잔액", "이 속도면 얼마가 남는지"), ("이번 달 손익", "확정 전표 기준"),
         ("받을 돈", "미수와 30일 초과"), ("낼 돈(30일)", "급여·정기지출·세금"), ("세금", "신고·납부 D-day")]
for i, (t, d) in enumerate(notes):
    x = Inches(0.62) + i * Inches(2.02)
    text(s, x, Inches(3.65), Inches(1.9), Inches(0.28), [(t, 11, True, BRAND)], PP_ALIGN.CENTER)
    text(s, x, Inches(3.98), Inches(1.9), Inches(0.5), [(d, 9.5, False, MUTED)], PP_ALIGN.CENTER, spacing=1.25)
rect(s, Inches(0.62), Inches(4.85), Inches(12.1), Inches(1.5), fill=SOFT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
text(s, Inches(1.0), Inches(5.12), Inches(11.4), Inches(1.0),
     [[("색이 곧 상태입니다.", 13, True, INK)],
      [("좋음 · 주의 · 위험 세 가지로만 칠합니다. 숫자를 읽지 않아도 오늘 회사가 어떤 상태인지 먼저 보입니다. "
        "각 칸을 누르면 그 숫자가 나온 화면으로 바로 이동합니다.", 11, False, MUTED)]], spacing=1.35)

# 13 · 04 오늘 챙길 것 (콜아웃형)
s = slide(); runner(s, "주요 기능"); footer(s, 13)
text(s, Inches(0.62), Inches(0.95), Inches(12.1), Inches(0.45),
     [("04. ", 20, True, BRAND), ("오늘 챙길 것을 ", 20, True, INK), ("먼저 알려 드립니다", 20, True, BRAND)])
text(s, Inches(0.62), Inches(1.52), Inches(12.1), Inches(0.3),
     [("화면을 뒤져 찾는 대신, 오늘 손대야 할 일이 근거와 함께 정리되어 있습니다.", 11, False, MUTED)])
pic(s, "03-briefing.png", Inches(0.62), Inches(2.05), Inches(8.6), Inches(4.5))
bx = Inches(9.45)
callouts = [("회사 숫자를 읽어 쓴 한 문단", "잔고·고정비·미수·세금 기한을 근거로 오늘의 요약을 만듭니다."),
            ("급한 순으로 정렬", "긴급·중요·권장으로 나누고, 금액과 경과일을 함께 적습니다."),
            ("바로 그 화면으로", "줄마다 '회수 관리 →'처럼 처리할 화면으로 이어집니다."),
            ("확정은 사람이", "AI 는 제안까지입니다. 독촉·이체·발송은 사람이 누릅니다.")]
for i, (t, d) in enumerate(callouts):
    y = Inches(2.15) + i * Inches(1.12)
    rect(s, bx, y, Inches(3.27), Inches(0.98), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    text(s, bx + Inches(0.2), y + Inches(0.14), Inches(2.9), Inches(0.28), [(t, 10.5, True, BRAND)])
    text(s, bx + Inches(0.2), y + Inches(0.45), Inches(2.9), Inches(0.5), [(d, 9, False, MUTED)], spacing=1.25)

# 14 · 05 분석
feature_rows(14, "05", "지금과 ", "30일 뒤를 함께 봅니다",
                 "과거를 정리해서 보여 주는 데서 멈추지 않고, 이 속도로 가면 어떻게 되는지까지 계산합니다.",
                 [("경영 흐름", "x/flow-chart.png",
                   "수입과 지출을 월별로 세워 흐름을 봅니다. 계정별로 펼쳐 어디서 늘고 줄었는지 확인합니다.", ""),
                  ("손익", "x/profit.png",
                   "확정 전표만으로 손익을 계산합니다. 아직 전표로 만들지 않은 자료는 몇 건인지 함께 알려 줍니다.", ""),
                  ("부가세 예상", "x/vat.png",
                   "이번 기의 매출·매입 자료로 낼 세금을 미리 셈합니다. 신고 기한은 신호 줄에서 D-day 로 보입니다.", "")])

# 15 · 06 프로젝트 (★ 항목 4개)
feature_sections(15, "06", "일마다 ", "맞는 화면을 드립니다",
                 "프로젝트를 만들면 하는 일에 맞는 첫 화면이 열립니다. 네 가지 예를 화면으로 보여 드립니다.",
                 [("템플릿에서 시작", "34-templates.png",
                   "부서가 아니라 일의 형태로 고릅니다. 고르면 칸까지 갖춰진 표가 만들어집니다.", ""),
                  ("회의 결정 → 할 일", "30-meeting-agenda.png",
                   "안건 아래 결정을 적고 버튼 하나로 담당·기한과 함께 넘깁니다.", "회의록이 실행으로"),
                  ("매출 파이프라인", "31-deal-card.png",
                   "검토→제안→계약→발행→입금을 점으로 봅니다. 잔금이 얼마 남았는지 카드에서 읽힙니다.", ""),
                  ("예산과 집행", "33-expense.png",
                   "남은 예산을 위에 붙박아 두고, 한 줄로 적습니다. 무엇에·어디에·얼마를 적고 Enter.", "")])

# 16 · 07 문서 흐름
feature_rows(16, "07", "견적에서 발행까지 ", "한 줄로 이어집니다",
                 "문서를 따로 만들지 않습니다. 앞 단계의 내용이 다음 문서로 그대로 넘어갑니다.",
                 [("견적서 · 계약서", "31-deal-card.png",
                   "매출 줄에서 바로 만듭니다. 앞 문서의 품목과 금액이 다음 문서로 그대로 넘어갑니다.", ""),
                  ("계약과 갱신", "x/pb-contract.png",
                   "체결한 계약은 만기 순으로 세워 둡니다. 지난 계약은 붉게 올라와 갱신을 놓치지 않습니다.", "종이 0장"),
                  ("세금계산서 발행", "x/tax-rows.png",
                   "계약 회차대로 청구 줄을 만들고 그대로 발행합니다. 발행 내역은 장부와 대조됩니다.", "")])

# 17 · 08 인사·협업
feature_rows(17, "08", "직원은 ", "마이페이지 하나로 끝냅니다",
                 "출퇴근부터 급여명세서까지 한 곳에서 처리하고, 관리자는 표 한 장으로 봅니다.",
                 [("구성원 관리", "x/employees.png",
                   "입사일·부서·직위·고용형태를 한 표로 봅니다. 재직·퇴직 인원과 인건비가 함께 셈해집니다.", ""),
                  ("근태", "x/attendance.png",
                   "출퇴근 기록이 쌓이고 지각·연장이 자동으로 계산됩니다. 급여 계산으로 이어집니다.", ""),
                  ("일정 · 협업", "x/schedule.png",
                   "일정과 할 일을 한 곳에서 봅니다. 공개 범위는 사람마다 정할 수 있습니다.", "")])

# 18 권한과 보안
s = slide(); runner(s, "주요 기능"); footer(s, 18)
head(s, [("권한은 ", 22, True, INK), ("메뉴 단위로, 보안은 데이터 단위로", 22, True, BRAND)],
     "누가 무엇을 볼 수 있는지 만들 때부터 정하고, 저장된 데이터에도 같은 규칙을 걸어 둡니다.")
pic(s, "x/perm.png", Inches(0.62), Inches(2.1), Inches(7.4), Inches(4.4))
sec = [("메뉴·세부 권한", "구성원마다 메뉴 접근과 세부 기능(금액 보기 등)을 따로 켭니다."),
       ("사용하게 될 메뉴 미리보기", "저장 전에 '○○님이 사용하게 될 메뉴'를 그대로 보여 줍니다."),
       ("행 단위 접근 통제(RLS)", "데이터베이스에서 회사·권한별로 막습니다. 화면만 가리지 않습니다."),
       ("감사 로그", "누가 언제 무엇을 바꿨는지 남습니다.")]
for i, (t, d) in enumerate(sec):
    y = Inches(2.2) + i * Inches(1.1)
    rect(s, Inches(8.3), y, Inches(4.42), Inches(0.95), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    text(s, Inches(8.55), y + Inches(0.13), Inches(4), Inches(0.28), [(t, 11, True, BRAND)])
    text(s, Inches(8.55), y + Inches(0.44), Inches(4), Inches(0.45), [(d, 9.5, False, MUTED)], spacing=1.25)

# 19 역할별 하루
s = slide(); runner(s, "주요 기능"); footer(s, 19)
head(s, [("역할마다 ", 22, True, INK), ("보는 화면이 다릅니다", 22, True, BRAND)],
     "같은 데이터를 쓰지만, 각자에게 필요한 것만 먼저 보입니다.")
roles = [("대표", "회사 신호 6칸 · 오늘 챙길 것 · 경영 흐름", "오늘 회사가 어떤 상태인지 30초 안에 확인합니다."),
         ("회계 담당", "수집 · 전표 · 세금계산서 · 부가세", "제안된 분류를 확인하고 확정합니다. 마감이 짧아집니다."),
         ("팀장", "프로젝트 · 결재 · 일정", "맡은 일의 진행과 결재를 한 화면에서 처리합니다."),
         ("직원", "마이페이지 · 근태 · 급여명세 · 요청", "출퇴근부터 증명서까지 스스로 처리합니다.")]
for i, (t, menus, d) in enumerate(roles):
    x = Inches(0.62) + i * Inches(3.09)
    rect(s, x, Inches(2.4), Inches(2.87), Inches(3.5), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, x, Inches(2.4), Inches(2.87), Inches(0.62), fill=BRAND, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    rect(s, x, Inches(2.85), Inches(2.87), Inches(0.17), fill=BRAND)
    text(s, x, Inches(2.55), Inches(2.87), Inches(0.32), [(t, 13, True, WHITE)], PP_ALIGN.CENTER)
    text(s, x + Inches(0.22), Inches(3.25), Inches(2.45), Inches(1.0), [(menus, 10, True, INK)], spacing=1.3)
    text(s, x + Inches(0.22), Inches(4.45), Inches(2.45), Inches(1.2), [(d, 10, False, MUTED)], spacing=1.3)

# 20 도입 효과
s = slide(); runner(s, "주요 기능"); footer(s, 20)
head(s, [("도입하면 ", 22, True, INK), ("이렇게 달라집니다", 22, True, BRAND)],
     "기능이 아니라 없어지는 일로 말씀드립니다.")
stats = [("7 → 1", "은행앱·엑셀·근태앱·문자로\n나뉘던 도구가 하나로"),
         ("0건", "통장·카드·세금계산서 수집에\n드는 손 입력"),
         ("D-5", "부가세·급여·계약 만기를\n미리 알려 드립니다"),
         ("0장", "계약과 결재에 드는\n종이와 출력")]
for i, (n, d) in enumerate(stats):
    x = Inches(0.62) + i * Inches(3.09)
    rect(s, x, Inches(2.55), Inches(2.87), Inches(2.6), fill=SOFT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(s, x, Inches(2.95), Inches(2.87), Inches(0.8), [(n, 34, True, BRAND)], PP_ALIGN.CENTER)
    text(s, x + Inches(0.2), Inches(3.95), Inches(2.47), Inches(1.0), [(d, 10.5, False, MUTED)], PP_ALIGN.CENTER, spacing=1.3)
rect(s, Inches(0.62), Inches(5.5), Inches(12.1), Inches(0.85), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
text(s, Inches(0.62), Inches(5.72), Inches(12.1), Inches(0.4),
     [("숫자를 모으는 시간이 사라지면, 그 시간은 판단에 쓰입니다.", 13, True, INK)], PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 21 챕터 Ⅳ
# ══════════════════════════════════════════════════════════
chapter("CHAPTER Ⅳ", "도입과 비용", "오늘 가입하고, 오늘부터 씁니다.", 21)

# 22 도입 절차
s = slide(); runner(s, "도입과 비용"); footer(s, 22)
head(s, [("도입은 ", 22, True, INK), ("네 단계로 끝납니다", 22, True, BRAND)],
     "설치가 필요 없습니다. 웹에서 가입하고 바로 씁니다.")
steps = [("가입", "이메일로 가입합니다.", "3분"),
         ("회사 개설", "사업자번호로 회사를 만듭니다. 이미 있으면 합류 요청을 보냅니다.", "5분"),
         ("연동", "공동인증서를 한 번 등록하면 은행·카드·홈택스 자료가 모입니다.", "10분"),
         ("구성원 초대", "메일로 초대하고, 권한을 메뉴별로 켜 줍니다.", "5분")]
for i, (t, d, time_) in enumerate(steps):
    x = Inches(0.62) + i * Inches(3.09)
    rect(s, x, Inches(2.5), Inches(2.87), Inches(2.9), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    c = rect(s, x + Inches(1.18), Inches(2.75), Inches(0.5), Inches(0.5), fill=BRAND, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(1.18), Inches(2.83), Inches(0.5), Inches(0.3), [(str(i + 1), 14, True, WHITE)], PP_ALIGN.CENTER)
    text(s, x, Inches(3.45), Inches(2.87), Inches(0.35), [(t, 13.5, True, INK)], PP_ALIGN.CENTER)
    text(s, x + Inches(0.25), Inches(3.9), Inches(2.37), Inches(1.0), [(d, 10, False, MUTED)], PP_ALIGN.CENTER, spacing=1.3)
    text(s, x, Inches(4.95), Inches(2.87), Inches(0.3), [(time_, 11, True, BRAND)], PP_ALIGN.CENTER)
    if i < 3:
        text(s, x + Inches(2.87), Inches(3.85), Inches(0.22), Inches(0.3), [("›", 18, True, DIM)], PP_ALIGN.CENTER)
rect(s, Inches(0.62), Inches(5.65), Inches(12.1), Inches(0.7), fill=SOFT, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
text(s, Inches(0.62), Inches(5.82), Inches(12.1), Inches(0.35),
     [("연동을 건너뛰어도 됩니다. 엑셀 업로드와 직접 입력으로 먼저 시작하고, 나중에 연결해도 자료는 그대로 이어집니다.", 11, False, MUTED)], PP_ALIGN.CENTER)

# 23 요금제
s = slide(); runner(s, "도입과 비용"); footer(s, 23)
head(s, [("무료로 시작해, ", 22, True, INK), ("회사에 맞게 자랍니다", 22, True, BRAND)],
     "기본 5명 포함 · VAT 별도 · 카드 등록 없이 무료로 시작합니다.")
plans = [("무료", "₩0", "카드 등록 없이 계속 무료",
          ["구성원 5명", "세금계산서 월 5건", "전자계약 월 5건", "통장·카드 3개 연결", "AI 참모 월 10만 토큰"], False),
         ("오너뷰", "₩39,000 / 월", "추가 1명 ₩5,000 · 가장 많이 쓰는 요금제",
          ["구성원 5명 + 추가 가능", "세금계산서 월 100건", "전자계약 무제한", "통장·카드 무제한 연결", "홈택스 자동 수집 · AI 브리핑"], True),
         ("울트라", "비용 협의", "회사 시스템에 맞춘 별도 UX 구축",
          ["구성원 무제한", "발행·연동 무제한", "화면·흐름 맞춤 설계", "기존 시스템 연동(ERP·그룹웨어)", "데이터 이관 · 전담 담당자"], False)]
for i, (t, price, sub, feats, hl) in enumerate(plans):
    x = Inches(0.62) + i * Inches(4.11)
    rect(s, x, Inches(2.35), Inches(3.88), Inches(4.05),
         fill=SOFT if hl else WHITE, line=BRAND if hl else LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x + Inches(0.3), Inches(2.62), Inches(3.3), Inches(0.35), [(t, 15, True, BRAND if hl else INK)])
    text(s, x + Inches(0.3), Inches(3.05), Inches(3.3), Inches(0.45), [(price, 20, True, INK)])
    text(s, x + Inches(0.3), Inches(3.62), Inches(3.3), Inches(0.4), [(sub, 9.5, False, DIM)], spacing=1.25)
    rect(s, x + Inches(0.3), Inches(4.12), Inches(3.28), Emu(9525), fill=LINE)
    for j, f in enumerate(feats):
        text(s, x + Inches(0.3), Inches(4.3) + j * Inches(0.38), Inches(3.3), Inches(0.3),
             [("· ", 10, True, BRAND), (f, 10, False, MUTED)])

# 24 시장과 사업 모델 (투자자)
s = slide(); runner(s, "투자자 참고"); footer(s, 24)
head(s, [("작은 회사의 운영 전체를 ", 22, True, INK), ("한 제품이 맡습니다", 22, True, BRAND)],
     "도구를 하나 더 파는 것이 아니라, 흩어진 도구를 대체하는 자리를 노립니다.")
cols = [("문제의 크기", "국내 법인 대부분이 은행앱·엑셀·세무 대행·근태앱·메신저를 따로 씁니다. "
                   "도구가 늘수록 숫자는 늦게 도착하고, 결정은 감에 기댑니다."),
        ("우리의 자리", "회계 SaaS 는 장부만, 그룹웨어는 결재만 봅니다. 오너뷰는 돈·일·사람·문서를 "
                   "한 데이터 위에 올려 '오늘 회사 상태'를 만들어 냅니다."),
        ("수익 구조", "월 구독(₩39,000~, 인원 과금)에 더해 울트라 맞춤 구축이 붙습니다. "
                   "구독은 예측 가능한 매출, 구축은 큰 계약과 이탈 방지를 맡습니다."),
        ("쌓이는 해자", "쓸수록 분류 규칙과 매칭 이력이 쌓여 제안이 정확해집니다. "
                   "장부·계약·인사 기록이 모두 안에 있어 옮기기 어려워집니다.")]
for i, (t, d) in enumerate(cols):
    x = Inches(0.62) + (i % 2) * Inches(6.16)
    y = Inches(2.4) + (i // 2) * Inches(2.15)
    rect(s, x, y, Inches(5.93), Inches(1.9), fill=WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(s, x + Inches(0.3), y + Inches(0.22), Inches(5.3), Inches(0.3), [(t, 13, True, BRAND)])
    text(s, x + Inches(0.3), y + Inches(0.65), Inches(5.3), Inches(1.1), [(d, 10.5, False, MUTED)], spacing=1.35)

# 25 로드맵 (투자자)
s = slide(); runner(s, "투자자 참고"); footer(s, 25)
head(s, [("지금 ", 22, True, INK), ("· 다음 · 그 다음", 22, True, BRAND)],
     "만든 것 위에 쌓습니다. 이미 도는 제품에서 시작합니다.")
phases = [("지금", "돌아가고 있습니다", ["은행·카드 자동 연동", "기장 자동화(분류·전표 제안)", "프로젝트·전자계약·결재", "인사·근태·급여", "AI 브리핑·참모"], True),
          ("다음", "정확도와 예측", ["매칭 정확도 향상(규칙 학습)", "자금 예측 고도화", "업종별 기본 설정", "모바일 사용성"], False),
          ("그 다음", "맞춤과 확장", ["업종별 울트라 패키지", "외부 시스템 연동(ERP·그룹웨어)", "세무 대리인 협업 기능", "해외(다국어) 검토"], False)]
for i, (t, sub, items_, hl) in enumerate(phases):
    x = Inches(0.62) + i * Inches(4.11)
    rect(s, x, Inches(2.4), Inches(3.88), Inches(3.9),
         fill=SOFT if hl else WHITE, line=BRAND if hl else LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    text(s, x + Inches(0.3), Inches(2.68), Inches(3.3), Inches(0.35), [(t, 15, True, BRAND if hl else INK)])
    text(s, x + Inches(0.3), Inches(3.1), Inches(3.3), Inches(0.3), [(sub, 10, False, DIM)])
    rect(s, x + Inches(0.3), Inches(3.5), Inches(3.28), Emu(9525), fill=LINE)
    for j, it in enumerate(items_):
        text(s, x + Inches(0.3), Inches(3.7) + j * Inches(0.42), Inches(3.3), Inches(0.35),
             [("· ", 10, True, BRAND), (it, 10, False, MUTED)])

# 26 마무리
s = slide()
gradient_bg(s, NAVY, BRAND, angle=45)
text(s, Inches(0.9), Inches(2.15), Inches(8.5), Inches(1.4),
     [[("회사 운영을 ", 28, True, WHITE), ("한 화면에서", 28, True, LAV)],
      [("시작해 보십시오.", 28, True, WHITE)]], spacing=1.25)
text(s, Inches(0.9), Inches(3.95), Inches(8), Inches(0.4),
     [("무료로 가입해 오늘 바로 쓰실 수 있습니다. 도입 상담과 맞춤 구축 문의를 받습니다.", 12, False, RGBColor(0xD8, 0xDC, 0xF5))])
info = [("서비스", "www.owner-view.com"), ("문의", "creative@mo-tive.com"), ("운영", "주식회사 모티브이노베이션")]
for i, (k, v) in enumerate(info):
    y = Inches(4.75) + i * Inches(0.55)
    text(s, Inches(0.9), y, Inches(1.5), Inches(0.35), [(k, 10.5, False, RGBColor(0xA5, 0xB4, 0xFC))])
    text(s, Inches(2.35), y, Inches(6), Inches(0.35), [(v, 12.5, True, WHITE)])
lg = rect(s, Inches(9.6), Inches(2.9), Inches(0.62), Inches(0.62), fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(9.6), Inches(3.02), Inches(0.62), Inches(0.4), [("O", 22, True, BRAND)], PP_ALIGN.CENTER)
text(s, Inches(10.35), Inches(3.0), Inches(3), Inches(0.5), [("OwnerView", 24, True, WHITE)])

os.makedirs("dist", exist_ok=True)
out = "dist/OwnerView_소개서_2026.pptx"
prs.save(out)
print("saved:", out)
