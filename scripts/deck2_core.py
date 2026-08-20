# -*- coding: utf-8 -*-
"""
OwnerView 서비스 소개서 v2 — 34장 (2026-08-20 사장님 확정 "진행")

읽는 사람: **도입을 검토하는 중소기업 대표 · 담당자 · 개인사업자**
  ① 구식 ERP가 무겁게 느껴지는 대표 ② 창업 초기라 아직 시스템이 없는 대표
  ③ 운영은 하는데 한눈에 안 보이는 대표 ④ 경리·총무 담당자 ⑤ 개인사업자

참고 덱에서 가져온 것 (사장님 지시로 역할을 나눠 분석):
  · 아드리엘 = **내실** — 기능 장표의 3단 밴드(제목/|부제| → 고객의 말 → 기능+캡처),
    연동 커넥터 그리드를 앞에 배치, 사용 예시 챕터, 화면이 분량의 70%
  · 마인드노크 = **외형** — 글꼴 굵기 3단계·중간 크기 배제, 캡처는 카드+그림자+캡션,
    3열 균등, 주색 3톤 + 강조 1색, 챕터에도 메시지 문장

데이터는 시연 회사 '오너뷰' — 모티브 실정보 노출 0.
실행: python scripts/deck2-build.py  →  dist/OwnerView_서비스소개서.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ── 색 ── (오너뷰 인디고 3톤 + 문제=빨강 + 효과=앰버)
BR   = RGBColor(0x1C, 0x3F, 0xAA)   # 주색
BR2  = RGBColor(0x3B, 0x5B, 0xDB)   # 밝은 주색
DEEP = RGBColor(0x0E, 0x1A, 0x44)   # 짙은 남색(표지·챕터)
BAND = RGBColor(0xEE, 0xF1, 0xF9)   # 밴드 배경
SOFT = RGBColor(0xF4, 0xF7, 0xFF)   # 카드 강조 배경
INK  = RGBColor(0x14, 0x18, 0x27)
MUT  = RGBColor(0x4A, 0x55, 0x70)
DIM  = RGBColor(0x93, 0x9B, 0xB1)
LINE = RGBColor(0xDF, 0xE3, 0xEE)
RED  = RGBColor(0xE0, 0x31, 0x31)
AMB  = RGBColor(0xB4, 0x53, 0x09)
WHITE= RGBColor(0xFF, 0xFF, 0xFF)
LAV  = RGBColor(0xC7, 0xD6, 0xFF)
BADR = RGBColor(0xFD, 0xF2, 0xF2)   # before 배경
BADL = RGBColor(0xF6, 0xD5, 0xD5)
GOODB= RGBColor(0xF2, 0xF6, 0xFF)
GOODL= RGBColor(0xCF, 0xDC, 0xFB)

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)
CAP = "cap/v2"
TOTAL = 34

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ══════════ 기본 도구 ══════════
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None, lw=0.75):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        for (t, size, bold, color) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
    return tb

# ══════════════════════════════════════════════════════════════
#  디자인 요소 (2026-08-20 사장님: "배경도 없고 밋밋하다 · 글자 크기·색·강조가 다 빠졌다")
#    참고: fs1.PNG(어두운 배경 + 배지 + 카드 3장), s133.PNG(강조어 색·형광 하이라이트)
# ══════════════════════════════════════════════════════════════
def _alpha(shape, pct):
    """도형에 투명도(%) — 배경 장식용. python-pptx 가 직접 지원하지 않아 XML 로 넣는다."""
    from pptx.oxml.ns import qn
    solid = shape.fill._xPr.find(qn("a:solidFill"))
    if solid is None:
        return shape
    clr = solid.find(qn("a:srgbClr"))
    if clr is None:
        return shape
    a = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - pct) * 1000))})
    clr.append(a)
    return shape


def bg_light(s):
    """밝은 장표 배경 — 옅은 하늘빛 그라데이션 + 큰 원 두 개. 흰 여백이 비어 보이지 않게."""
    grad(s, RGBColor(0xF8, 0xFA, 0xFF), RGBColor(0xE9, 0xEE, 0xFB), 45)
    c1 = rect(s, W - Inches(3.2), -Inches(2.1), Inches(6.4), Inches(6.4), fill=BR2, shape=MSO_SHAPE.OVAL)
    _alpha(c1, 93)
    c2 = rect(s, -Inches(2.4), H - Inches(2.6), Inches(5.2), Inches(5.2), fill=BR, shape=MSO_SHAPE.OVAL)
    _alpha(c2, 95)
    bar = rect(s, 0, 0, W, Inches(0.075), fill=BR2)
    return s


def bg_dark(s):
    """어두운 장표 배경 — 딥 인디고 그라데이션 + 대각선 띠. 표지·챕터·선언 장표용."""
    grad(s, RGBColor(0x0C, 0x14, 0x36), RGBColor(0x2B, 0x3F, 0xA8), 35)
    d1 = rect(s, Inches(7.6), -Inches(1.2), Inches(3.1), Inches(10), fill=RGBColor(0x6C, 0x8C, 0xFF),
              shape=MSO_SHAPE.PARALLELOGRAM)
    _alpha(d1, 90)
    d2 = rect(s, Inches(10.4), -Inches(1.6), Inches(2.2), Inches(10), fill=RGBColor(0x9F, 0xB4, 0xFF),
              shape=MSO_SHAPE.PARALLELOGRAM)
    _alpha(d2, 93)
    c = rect(s, -Inches(1.6), H - Inches(3.0), Inches(5.6), Inches(5.6), fill=RGBColor(0x8F, 0xA5, 0xF5),
             shape=MSO_SHAPE.OVAL)
    _alpha(c, 94)
    return s


def pill(s, cx, y, label, w=Inches(2.6), dark=False):
    """제목 위 알약 배지 — 참고 덱의 '풍부한 경험과 노하우' 자리"""
    x = cx - w // 2
    h = Inches(0.34)
    sp = rect(s, x, y, w, h, fill=WHITE if dark else BR, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, x, y + Inches(0.055), w, Inches(0.26),
         [(label, 10.5, True, BR if dark else WHITE)], PP_ALIGN.CENTER)
    return sp


def mark(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=None, hl=None):
    """강조어에 형광 배경을 깔고 글자를 얹는다 (s133.PNG 의 '사업 목표를 실현' 표현)"""
    if hl:
        hx, hw = hl
        rect(s, hx, y + Inches(0.04), hw, h - Inches(0.06), fill=RGBColor(0xE4, 0xEA, 0xFE))
    return text(s, x, y, w, h, runs, align, spacing=spacing)


ICON_FONT = "Segoe MDL2 Assets"
#   Windows 기본 아이콘 폰트 — PowerPoint·PDF 모두에서 깨지지 않는다.
ICON = {
    "bank": "", "card": "", "chart": "", "people": "",
    "doc": "", "lock": "", "calendar": "", "gear": "",
    "money": "", "check": "", "search": "", "flag": "",
    "clock": "", "cloud": "", "mail": "", "star": "",
    "shield": "", "bulb": "", "list": "", "link": "",
}


def icon_badge(s, cx, cy, key, d=Inches(0.86), no=None, tone=None):
    """원형 아이콘 배지 (+ 우상단 번호) — 참고 덱의 프로세스 아이콘 표현"""
    x, y = cx - d // 2, cy - d // 2
    base = tone if tone is not None else RGBColor(0xE4, 0xEA, 0xFE)
    rect(s, x, y, d, d, fill=base, shape=MSO_SHAPE.OVAL)
    tb = s.shapes.add_textbox(x, y + int(d * 0.22), d, int(d * 0.6))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ICON.get(key, ICON["check"])
    r.font.size = Pt(int(d / 914400 * 26)); r.font.bold = False
    r.font.color.rgb = BR; r.font.name = ICON_FONT
    if no is not None:
        nd = Inches(0.3)
        nx, ny = x + d - nd + Inches(0.02), y - Inches(0.02)
        rect(s, nx, ny, nd, nd, fill=BR, shape=MSO_SHAPE.OVAL)
        text(s, nx, ny + Inches(0.045), nd, Inches(0.22), [(str(no), 9, True, WHITE)], PP_ALIGN.CENTER)
    return s


def check_list(s, x, y, w, items, hl=False):
    """✔ 체크 목록 박스 — 참고 덱 하단의 항목 상자"""
    h = Inches(0.34) * len(items) + Inches(0.3)
    rect(s, x, y, w, h, fill=SOFT if hl else WHITE, line=BR if hl else LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    for i, it in enumerate(items):
        yy = y + Inches(0.15) + i * Inches(0.34)
        tb = s.shapes.add_textbox(x + Inches(0.18), yy, Inches(0.22), Inches(0.24))
        tf = tb.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        r = tf.paragraphs[0].add_run(); r.text = ICON["check"]
        r.font.size = Pt(8); r.font.color.rgb = BR; r.font.name = ICON_FONT
        text(s, x + Inches(0.44), yy, w - Inches(0.6), Inches(0.26), [(it, 9.5, False, MUT)])
    return h


def avatar(s, cx, cy, d=Inches(1.05), tone=None, tie=None):
    """플랫 아바타 — 도형 조합(머리·몸·옷깃). 참고 덱의 캐릭터 자리를 대신한다."""
    body = tone if tone is not None else RGBColor(0x4F, 0x6F, 0xF0)
    skin = RGBColor(0xFF, 0xD9, 0xBE)
    hair = RGBColor(0x2A, 0x2F, 0x45)
    # 몸(어깨)
    bw, bh = int(d * 1.15), int(d * 0.62)
    rect(s, cx - bw // 2, cy + int(d * 0.16), bw, bh, fill=body, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.45)
    # 옷깃
    if tie is not None:
        rect(s, cx - int(d * 0.07), cy + int(d * 0.2), int(d * 0.14), int(d * 0.34), fill=tie)
    # 머리
    hd = int(d * 0.58)
    hy = cy - int(d * 0.46)
    rect(s, cx - hd // 2, hy, hd, hd, fill=skin, shape=MSO_SHAPE.OVAL)
    # 머리카락 — 정수리만 덮는다(얼굴을 가리지 않게)
    rect(s, cx - int(hd * 0.54), hy - int(hd * 0.1), int(hd * 1.08), int(hd * 0.46),
         fill=hair, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    return s


def bubble(s, x, y, w, h, runs, fill=None, tail="left", align=PP_ALIGN.LEFT, spacing=1.3, pad=Inches(0.3)):
    """말풍선 — 둥근 상자 + 꼬리. Q/A·고객의 말에 쓴다."""
    bg = fill if fill is not None else RGBColor(0xE8, 0xED, 0xFD)
    rect(s, x, y, w, h, fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    t = Inches(0.16)
    if tail == "left":
        tr = rect(s, x + Inches(0.5), y + h - Inches(0.02), t, t, fill=bg, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        tr.rotation = 180
    elif tail == "topleft":
        tr = rect(s, x + Inches(0.5), y - t + Inches(0.02), t, t, fill=bg, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
    text(s, x + pad, y + Inches(0.16), w - pad * 2, h - Inches(0.3), runs, align, spacing=spacing)
    return s


def link_text(s, x, y, w, h, label, url, size=12, bold=True, color=None, align=PP_ALIGN.LEFT):
    """하이퍼링크가 걸린 글자 — 발표·PDF 어느 쪽에서 눌러도 그 페이지로 간다 (2026-08-20 사장님 지시)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT
    r.font.color.rgb = color if color is not None else BR
    r.hyperlink.address = url
    return tb


def link_shape(s, shape, url):
    """도형(버튼) 자체에 링크를 건다"""
    from pptx.oxml.ns import qn
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    rId = s.part.relate_to(url, RT.HYPERLINK, is_external=True)
    cNvPr = shape._element._nvXxPr.cNvPr
    hl = cNvPr.makeelement(qn("a:hlinkClick"), {qn("r:id"): rId})
    cNvPr.append(hl)
    return shape


def grad(s, c1, c2, angle=45):
    sp = rect(s, 0, 0, W, H, fill=c1)
    f = sp.fill; f.gradient(); f.gradient_angle = angle
    st = f.gradient_stops
    st[0].color.rgb = c1; st[0].position = 0.0
    st[1].color.rgb = c2; st[1].position = 1.0
    return sp

def pic(s, name, x, y, maxw, maxh, border=True, align="center"):
    """비율 유지 — 영역 안에 가운데 배치"""
    p = os.path.join(CAP, name)
    if not os.path.exists(p):
        rect(s, x, y, maxw, maxh, fill=BAND, line=LINE)
        text(s, x, y + maxh // 2, maxw, Inches(0.3), [(f"[{name}]", 9, False, DIM)], PP_ALIGN.CENTER)
        return
    iw, ih = Image.open(p).size
    r = min(maxw / iw, maxh / ih)
    w_, h_ = int(iw * r), int(ih * r)
    px = x if align == "left" else x + (maxw - w_) // 2
    py = y + (maxh - h_) // 2
    s.shapes.add_picture(p, px, py, w_, h_)
    if border:
        b = rect(s, px, py, w_, h_, fill=None, line=LINE); b.line.width = Pt(0.75)

def pic_top(s, name, x, y, w, h):
    """폭을 맞추고 위에서부터 잘라 넣는다 — 화면 전체를 크게 보여줄 때"""
    p = os.path.join(CAP, name)
    if not os.path.exists(p):
        rect(s, x, y, w, h, fill=BAND, line=LINE); return
    iw, ih = Image.open(p).size
    ar, tr = iw / ih, w / h
    im = s.shapes.add_picture(p, x, y, w, h)
    if ar > tr:
        cut = 1 - tr / ar; im.crop_right = cut
    elif ar < tr:
        cut = 1 - ar / tr; im.crop_bottom = cut
    b = rect(s, x, y, w, h, fill=None, line=LINE); b.line.width = Pt(0.75)

def foot(s, n):
    text(s, Inches(0.55), Inches(7.02), Inches(3), Inches(0.22), [("OwnerView", 8, True, BR)])
    text(s, Inches(9.8), Inches(7.02), Inches(3), Inches(0.22),
         [(f"{n:02d} / {TOTAL}", 8, False, DIM)], PP_ALIGN.RIGHT)


# ══════════ 장표 틀 ══════════
def cover():
    s = slide(); bg_dark(s)
    rect(s, 0, 0, W, Inches(0.06), fill=LAV)
    text(s, Inches(0.85), Inches(0.75), Inches(8), Inches(0.3),
         [("COMPANY OPERATION SOFTWARE", 10.5, True, LAV)])
    text(s, Inches(0.85), Inches(2.4), Inches(7.2), Inches(2),
         [[("은행·카드·세금·계약·직원 —", 28, True, WHITE)],
          [("회사 운영의 모든 것을 한 화면에", 28, True, LAV)]], spacing=1.28)
    text(s, Inches(0.85), Inches(5.3), Inches(5), Inches(0.4), [("◈ OwnerView", 19, True, WHITE)])
    text(s, Inches(0.85), Inches(5.95), Inches(4.05), Inches(0.35),
         [("작은 회사를 위한 운영 소프트웨어 · ", 11.5, False, RGBColor(0xD5, 0xDC, 0xF5))])
    #   표지 주소 — 글자는 흰색으로 두고 링크는 투명 도형에 건다
    #   (하이퍼링크 run 은 테마색(파랑·밑줄)이 강제돼 표지에서 튄다)
    text(s, Inches(4.35), Inches(5.95), Inches(3.4), Inches(0.35),
         [("www.owner-view.com", 11.5, True, RGBColor(0xFF, 0xFF, 0xFF))])
    hot = rect(s, Inches(4.3), Inches(5.92), Inches(2.5), Inches(0.32), fill=None)
    link_shape(s, hot, "https://www.owner-view.com")
    p = os.path.join(CAP, "use-dash.png")
    if os.path.exists(p):
        iw, ih = Image.open(p).size
        w_ = Inches(6.5); h_ = int(w_ * ih / iw)
        s.shapes.add_picture(p, Inches(7.1), Inches(1.5), w_, h_)
    return s

def chapter(no, title, sub, page):
    s = slide(); bg_dark(s)
    rect(s, 0, 0, Inches(0.15), H, fill=BR2)
    text(s, Inches(1.1), Inches(2.8), Inches(8), Inches(0.3), [(no, 11.5, True, LAV)])
    text(s, Inches(1.1), Inches(3.2), Inches(9), Inches(0.9), [(title, 30, True, WHITE)])
    text(s, Inches(1.1), Inches(4.2), Inches(9), Inches(0.4), [(sub, 11.5, False, RGBColor(0xC7, 0xCB, 0xE0))])
    text(s, Inches(9.8), Inches(7.02), Inches(3), Inches(0.22),
         [(f"{page:02d} / {TOTAL}", 8, False, RGBColor(0x8A, 0x90, 0xB5))], PP_ALIGN.RIGHT)
    return s

def head_band(s, title, pipe, badge=None):
    """상단 밴드 — 배지 + 제목 + |부제|"""
    b = rect(s, 0, 0, W, Inches(1.62), fill=BAND)
    f = b.fill; f.gradient(); f.gradient_angle = 90
    st = f.gradient_stops
    st[0].color.rgb = RGBColor(0xE6, 0xEC, 0xFA); st[0].position = 0.0
    st[1].color.rgb = RGBColor(0xF4, 0xF7, 0xFF); st[1].position = 1.0
    rect(s, 0, Inches(1.62), W, Inches(0.035), fill=BR2)
    if badge:
        pill(s, W // 2, Inches(0.18), badge, w=Inches(2.5))
    text(s, Inches(0.6), Inches(0.62), Inches(12.13), Inches(0.5),
         [(title, 23, True, INK)], PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(1.2), Inches(12.13), Inches(0.3),
         [("| ", 12.5, True, BR), (pipe, 12.5, True, BR), (" |", 12.5, True, BR)], PP_ALIGN.CENTER)

def three_band(page, title, pipe, voices, feats):
    """기능 장표 — 밴드(제목/|부제|) / 고객의 말 2 / 왼쪽 기능 설명 3 + 오른쪽 **화면 전체 1장**

       2026-08-20 사장님: "이미지 전부 다 바꿔줘 · 전체화면이 다 나오게".
       조각을 3개 늘어놓으면 각각이 작아진다 — 대표 화면 하나를 크게 두고 설명은 왼쪽에 쌓는다.
    """
    s = slide(); bg_light(s); head_band(s, title, pipe, badge="주요 기능")
    y = Inches(2.0)
    for (a, b, c) in voices:
        runs = [("“" + a, 13, True, INK), (b, 13, True, RED), (c + "”", 13, True, INK)]
        text(s, Inches(1.0), y, Inches(11.3), Inches(0.4), runs, PP_ALIGN.CENTER)
        y += Inches(0.5)
    # 왼쪽 — 기능 설명 3개(아이콘 배지 포함)
    ICONS = ["check", "chart", "bulb"]
    ty = Inches(3.15)
    for i, (nm, desc, _img) in enumerate(feats):
        rect(s, Inches(0.6), ty, Inches(4.35), Inches(1.2), fill=WHITE, line=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
        rect(s, Inches(0.6), ty, Inches(0.09), Inches(1.2), fill=BR)
        icon_badge(s, Inches(1.12), ty + Inches(0.6), ICONS[i % 3], d=Inches(0.5), no=i + 1)
        text(s, Inches(1.52), ty + Inches(0.18), Inches(3.3), Inches(0.3), [(nm, 12.5, True, BR)])
        text(s, Inches(1.52), ty + Inches(0.52), Inches(3.3), Inches(0.55),
             [(desc.replace(chr(10), " "), 9.5, False, MUT)], spacing=1.25)
        ty += Inches(1.34)

    # 오른쪽 — 대표 화면 전체(자르지 않는다)
    hero = feats[0][2]
    rect(s, Inches(5.2), Inches(3.15), Inches(7.53), Inches(3.72), fill=WHITE, line=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
    pic(s, hero, Inches(5.32), Inches(3.27), Inches(7.29), Inches(3.48), border=False)
    foot(s, page)
    return s


def use_slide(page, title, pipes, img):
    """사용 예시 — 제목 좌상단 + |부제 2줄| + 화면 크게"""
    s = slide(); bg_light(s)
    rect(s, Inches(0.55), Inches(0.42), Inches(1.55), Inches(0.32), fill=BR,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(s, Inches(0.55), Inches(0.47), Inches(1.55), Inches(0.24),
         [("사용 예시", 10, True, WHITE)], PP_ALIGN.CENTER)
    text(s, Inches(2.25), Inches(0.38), Inches(10.5), Inches(0.45), [(title, 19, True, INK)])
    y = Inches(0.98)
    for t in pipes:
        text(s, Inches(0.55), y, Inches(12.2), Inches(0.28),
             [("| ", 11, True, RED), (t, 11, False, RGBColor(0x3A, 0x43, 0x56))])
        y += Inches(0.3)
    #   화면 전체가 한눈에 들어와야 한다 — 자르지 않고 비율을 지켜 넣는다 (2026-08-20 사장님 지시)
    box_y = Inches(1.62)
    pic(s, img, Inches(0.55), box_y, Inches(12.23), H - box_y - Inches(0.38))
    foot(s, page)
    return s

def plain(page, title, lead=None, band=False):
    s = slide(); bg_light(s)
    if band:
        head_band(s, title, lead or "")
    else:
        text(s, Inches(0.6), Inches(0.55), Inches(12.13), Inches(0.5), [(title, 23, True, INK)], PP_ALIGN.CENTER)
        if lead:
            text(s, Inches(0.6), Inches(1.15), Inches(12.13), Inches(0.35), [(lead, 11.5, False, MUT)], PP_ALIGN.CENTER)
    foot(s, page)
    return s
